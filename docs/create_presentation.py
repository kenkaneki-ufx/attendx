"""
AttendX Project PowerPoint Presentation Generator
Creates a professional 12-slide presentation for B.Tech CSE Mini Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def create_attendx_presentation():
    """Create the AttendX PowerPoint presentation."""
    
    # Create presentation with widescreen layout
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PRIMARY_BLUE = RGBColor(13, 110, 253)  # Bootstrap primary blue
    DARK_BLUE = RGBColor(8, 66, 152)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(248, 249, 250)
    DARK_GRAY = RGBColor(33, 37, 41)
    SUCCESS_GREEN = RGBColor(25, 135, 84)
    WARNING_ORANGE = RGBColor(255, 193, 7)
    
    def add_background(slide, color):
        """Add colored background to slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_shape(slide, left, top, width, height, color, opacity=100):
        """Add a colored shape to slide."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
        """Add a text box with formatted text."""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        
        return text_box
    
    def add_bullet_points(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
        """Add bullet points to slide."""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.space_after = Pt(8)
        
        return text_box
    
    def add_slide_number(slide, num, total=12):
        """Add slide number to bottom right."""
        text_box = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"{num} / {total}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.RIGHT
    
    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide1, PRIMARY_BLUE)
    
    # Title
    add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "AttendX", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Subtitle
    add_text_box(slide1, Inches(1), Inches(3), Inches(11), Inches(1),
                 "QR Code Attendance Management System", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Tagline
    add_text_box(slide1, Inches(1), Inches(4), Inches(11), Inches(0.8),
                 "Scan • Verify • Attend", font_size=24, color=WARNING_ORANGE, alignment=PP_ALIGN.CENTER)
    
    # Project details
    details = [
        "B.Tech CSE Mini Project",
        "Department of Computer Science & Engineering",
        "Academic Year: 2025-2026"
    ]
    add_bullet_points(slide1, Inches(3), Inches(5.2), Inches(7), Inches(1.5),
                     details, font_size=18, color=WHITE)
    
    add_slide_number(slide1, 1)
    
    # ==================== SLIDE 2: INTRODUCTION ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide2, WHITE)
    
    # Header bar
    add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide2, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Introduction", font_size=36, bold=True, color=WHITE)
    
    # Content
    intro_text = """AttendX is a modern, secure, and efficient QR Code-based Attendance Management System 
designed specifically for educational institutions. It revolutionizes the traditional attendance 
tracking process by eliminating manual roll calls and paper-based registers."""
    
    add_text_box(slide2, Inches(0.5), Inches(1.5), Inches(12), Inches(1.5),
                 intro_text, font_size=18, color=DARK_GRAY)
    
    # Key highlights
    highlights = [
        "Faculty members generate dynamic, time-limited QR codes for each lecture",
        "Students simply scan the QR code using any device camera (no app required)",
        "Real-time attendance tracking with live dashboard updates",
        "Automatic absence marking when lectures end",
        "Comprehensive analytics and report generation",
        "Role-based access control for security"
    ]
    add_bullet_points(slide2, Inches(0.5), Inches(3.2), Inches(12), Inches(4),
                     highlights, font_size=16, color=DARK_GRAY)
    
    add_slide_number(slide2, 2)
    
    # ==================== SLIDE 3: PROBLEM STATEMENT ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide3, WHITE)
    
    # Header bar
    add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Problem Statement", font_size=36, bold=True, color=WHITE)
    
    # Problem description
    problem_text = """Traditional attendance systems in educational institutions face numerous challenges 
that affect efficiency, accuracy, and security:"""
    add_text_box(slide3, Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                 problem_text, font_size=18, color=DARK_GRAY)
    
    # Problems list
    problems = [
        "Time-consuming manual roll calls waste valuable lecture time",
        "Paper-based registers are prone to errors and difficult to maintain",
        "Proxy attendance (buddy punching) is easy to commit and hard to detect",
        "Generating accurate attendance reports requires hours of manual work",
        "No real-time visibility into attendance patterns for faculty",
        "Students can mark attendance without actually being present",
        "Difficulty in tracking attendance across multiple subjects and semesters"
    ]
    add_bullet_points(slide3, Inches(0.5), Inches(2.8), Inches(12), Inches(4.5),
                     problems, font_size=16, color=DARK_GRAY)
    
    add_slide_number(slide3, 3)
    
    # ==================== SLIDE 4: OBJECTIVES ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide4, WHITE)
    
    # Header bar
    add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide4, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Objectives & Goals", font_size=36, bold=True, color=WHITE)
    
    # Primary objectives
    add_text_box(slide4, Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
                 "Primary Objectives:", font_size=22, bold=True, color=DARK_BLUE)
    
    primary = [
        "Develop a secure, web-based attendance system using QR code technology",
        "Enable contactless attendance marking without requiring student mobile apps",
        "Implement real-time attendance tracking with live dashboard updates",
        "Automate attendance report generation in multiple formats (PDF, Excel, CSV)",
        "Provide analytics and insights for faculty and administrators"
    ]
    add_bullet_points(slide4, Inches(0.5), Inches(2.2), Inches(12), Inches(2.5),
                     primary, font_size=16, color=DARK_GRAY)
    
    # Secondary objectives
    add_text_box(slide4, Inches(0.5), Inches(4.8), Inches(12), Inches(0.6),
                 "Secondary Goals:", font_size=22, bold=True, color=DARK_BLUE)
    
    secondary = [
        "Prevent proxy attendance and ensure attendance integrity",
        "Create a scalable architecture ready for production deployment",
        "Follow software engineering best practices and clean code principles",
        "Build a professional portfolio-quality application"
    ]
    add_bullet_points(slide4, Inches(0.5), Inches(5.5), Inches(12), Inches(2),
                     secondary, font_size=16, color=DARK_GRAY)
    
    add_slide_number(slide4, 4)
    
    # ==================== SLIDE 5: TECHNOLOGIES USED ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide5, WHITE)
    
    # Header bar
    add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Technologies Used", font_size=36, bold=True, color=WHITE)
    
    # Technology categories
    # Backend
    add_shape(slide5, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_text_box(slide5, Inches(0.7), Inches(1.6), Inches(3.4), Inches(0.5),
                 "Backend", font_size=20, bold=True, color=PRIMARY_BLUE)
    backend_tech = ["Python 3.14", "Django 5.2 LTS", "Django ORM", "SQLite / MySQL"]
    add_bullet_points(slide5, Inches(0.7), Inches(2.2), Inches(3.4), Inches(1.8),
                     backend_tech, font_size=14, color=DARK_GRAY)
    
    # Frontend
    add_shape(slide5, Inches(4.8), Inches(1.5), Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_text_box(slide5, Inches(5), Inches(1.6), Inches(3.4), Inches(0.5),
                 "Frontend", font_size=20, bold=True, color=PRIMARY_BLUE)
    frontend_tech = ["HTML5 / CSS3", "Bootstrap 5.3", "JavaScript ES6", "AJAX / Fetch API"]
    add_bullet_points(slide5, Inches(5), Inches(2.2), Inches(3.4), Inches(1.8),
                     frontend_tech, font_size=14, color=DARK_GRAY)
    
    # Libraries
    add_shape(slide5, Inches(9.1), Inches(1.5), Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_text_box(slide5, Inches(9.3), Inches(1.6), Inches(3.4), Inches(0.5),
                 "Libraries", font_size=20, bold=True, color=PRIMARY_BLUE)
    libraries = ["qrcode (QR Generation)", "Pillow (Image Processing)", "Chart.js (Analytics)", "ReportLab (PDF)"]
    add_bullet_points(slide5, Inches(9.3), Inches(2.2), Inches(3.4), Inches(1.8),
                     libraries, font_size=14, color=DARK_GRAY)
    
    # Security & Tools
    add_shape(slide5, Inches(0.5), Inches(4.3), Inches(6), Inches(2.8), LIGHT_GRAY)
    add_text_box(slide5, Inches(0.7), Inches(4.4), Inches(5.6), Inches(0.5),
                 "Security & Tools", font_size=20, bold=True, color=PRIMARY_BLUE)
    security = [
        "CSRF Protection (Django Built-in)",
        "Session-based Authentication",
        "Password Hashing (PBKDF2)",
        "Rate Limiting",
        "Secure Token Generation"
    ]
    add_bullet_points(slide5, Inches(0.7), Inches(5), Inches(5.6), Inches(2),
                     security, font_size=14, color=DARK_GRAY)
    
    # Development Tools
    add_shape(slide5, Inches(6.8), Inches(4.3), Inches(6.1), Inches(2.8), LIGHT_GRAY)
    add_text_box(slide5, Inches(7), Inches(4.4), Inches(5.7), Inches(0.5),
                 "Development Tools", font_size=20, bold=True, color=PRIMARY_BLUE)
    tools = ["Git (Version Control)", "VS Code (IDE)", "Pytest (Testing)", "Ruff (Linting)", "Docker (Optional)"]
    add_bullet_points(slide5, Inches(7), Inches(5), Inches(5.7), Inches(2),
                     tools, font_size=14, color=DARK_GRAY)
    
    add_slide_number(slide5, 5)
    
    # ==================== SLIDE 6: SYSTEM ARCHITECTURE ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide6, WHITE)
    
    # Header bar
    add_shape(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide6, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "System Architecture", font_size=36, bold=True, color=WHITE)
    
    # Architecture layers
    # Presentation Layer
    add_shape(slide6, Inches(1), Inches(1.5), Inches(11.333), Inches(1.5), RGBColor(227, 242, 253))
    add_text_box(slide6, Inches(1.2), Inches(1.6), Inches(3), Inches(0.5),
                 "Presentation Layer", font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide6, Inches(1.2), Inches(2.1), Inches(10.5), Inches(0.8),
                 "HTML5 | CSS3 | Bootstrap 5 | JavaScript | AJAX | Templates", font_size=14, color=DARK_GRAY)
    
    # Arrow down
    add_text_box(slide6, Inches(6), Inches(3), Inches(1.5), Inches(0.5),
                 "▼", font_size=24, color=PRIMARY_BLUE, alignment=PP_ALIGN.CENTER)
    
    # Application Layer
    add_shape(slide6, Inches(1), Inches(3.5), Inches(11.333), Inches(1.5), RGBColor(209, 236, 241))
    add_text_box(slide6, Inches(1.2), Inches(3.6), Inches(3), Inches(0.5),
                 "Application Layer", font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide6, Inches(1.2), Inches(4.1), Inches(10.5), Inches(0.8),
                 "Django Views | URL Routing | Middleware | Service Layer | Selectors", font_size=14, color=DARK_GRAY)
    
    # Arrow down
    add_text_box(slide6, Inches(6), Inches(5), Inches(1.5), Inches(0.5),
                 "▼", font_size=24, color=PRIMARY_BLUE, alignment=PP_ALIGN.CENTER)
    
    # Data Layer
    add_shape(slide6, Inches(1), Inches(5.5), Inches(11.333), Inches(1.5), RGBColor(178, 223, 219))
    add_text_box(slide6, Inches(1.2), Inches(5.6), Inches(3), Inches(0.5),
                 "Data Layer", font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide6, Inches(1.2), Inches(6.1), Inches(10.5), Inches(0.8),
                 "Django ORM | Models | Migrations | SQLite (Dev) / MySQL (Prod)", font_size=14, color=DARK_GRAY)
    
    add_slide_number(slide6, 6)
    
    # ==================== SLIDE 7: WORKING (FLOWCHART) ====================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide7, WHITE)
    
    # Header bar
    add_shape(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide7, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Working Process", font_size=36, bold=True, color=WHITE)
    
    # Flowchart steps
    steps = [
        ("1. Faculty Login", Inches(0.5), PRIMARY_BLUE),
        ("2. Start Lecture", Inches(2.5), DARK_BLUE),
        ("3. Generate QR", Inches(4.5), PRIMARY_BLUE),
        ("4. Display QR", Inches(6.5), DARK_BLUE),
        ("5. Student Scans", Inches(8.5), PRIMARY_BLUE),
        ("6. Verify & Record", Inches(10.5), DARK_BLUE),
    ]
    
    for step_text, left, color in steps:
        shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = step_text
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow between steps (except last)
        if left < Inches(10.5):
            add_text_box(slide7, left + Inches(2), Inches(2), Inches(0.5), Inches(0.5),
                        "→", font_size=20, color=PRIMARY_BLUE, alignment=PP_ALIGN.CENTER)
    
    # Detailed workflow
    add_text_box(slide7, Inches(0.5), Inches(3.2), Inches(12), Inches(0.6),
                 "Detailed Workflow:", font_size=20, bold=True, color=DARK_BLUE)
    
    workflow = [
        "Faculty logs in securely and navigates to the dashboard",
        "Faculty selects subject, section, and starts a new lecture",
        "System generates a unique, time-limited QR code (60-second expiry)",
        "QR code is displayed on the projector/screen for students to scan",
        "Students scan QR code using any device camera (no app required)",
        "System validates QR token, checks expiration, and records attendance",
        "Live attendance counter updates in real-time on faculty dashboard",
        "When lecture ends, unmarked students are automatically marked absent"
    ]
    add_bullet_points(slide7, Inches(0.5), Inches(3.9), Inches(12), Inches(3.5),
                     workflow, font_size=15, color=DARK_GRAY)
    
    add_slide_number(slide7, 7)
    
    # ==================== SLIDE 8: WORKING METHODOLOGY ====================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide8, WHITE)
    
    # Header bar
    add_shape(slide8, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide8, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Working Methodology", font_size=36, bold=True, color=WHITE)
    
    # Methodology description
    add_text_box(slide8, Inches(0.5), Inches(1.5), Inches(12), Inches(0.8),
                 "AttendX follows a service-layer architecture pattern with clear separation of concerns:",
                 font_size=18, color=DARK_GRAY)
    
    # Methodology boxes
    methods = [
        ("Service Layer", "Business logic encapsulation in services.py files. Handles QR generation, attendance processing, and report generation."),
        ("Selector Layer", "Complex database queries centralized in selectors.py. Optimizes queries with select_related and prefetch_related."),
        ("Template Pattern", "Reusable HTML templates with Bootstrap 5. Dynamic content rendering with Django template language."),
        ("Security First", "Cryptographic tokens, CSRF protection, rate limiting, and secure session management throughout."),
    ]
    
    for i, (title, desc) in enumerate(methods):
        left = Inches(0.5) if i % 2 == 0 else Inches(6.8)
        top = Inches(2.5) if i < 2 else Inches(4.5)
        
        add_shape(slide8, left, top, Inches(5.8), Inches(1.8), LIGHT_GRAY)
        add_text_box(slide8, left + Inches(0.2), top + Inches(0.1), Inches(5.4), Inches(0.5),
                    title, font_size=18, bold=True, color=PRIMARY_BLUE)
        add_text_box(slide8, left + Inches(0.2), top + Inches(0.6), Inches(5.4), Inches(1),
                    desc, font_size=14, color=DARK_GRAY)
    
    add_slide_number(slide8, 8)
    
    # ==================== SLIDE 9: MODULES AND PAGES ====================
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide9, WHITE)
    
    # Header bar
    add_shape(slide9, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide9, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Modules & Pages", font_size=36, bold=True, color=WHITE)
    
    # Modules
    modules = [
        ("Authentication", "Login, Logout, Profile, Password Reset"),
        ("Dashboard", "Faculty Dashboard, Quick Stats, Recent Activity"),
        ("Lectures", "Start Lecture, Active Lecture, End Lecture"),
        ("QR Codes", "Generate QR, Regenerate, QR History"),
        ("Attendance", "Live Tracking, History, Auto Absent"),
        ("Reports", "Daily, Weekly, Monthly, Subject-wise Reports"),
        ("Analytics", "Charts, Trends, Performance Metrics"),
        ("Administration", "Faculty/Student/Subject Management"),
    ]
    
    for i, (module, pages) in enumerate(modules):
        left = Inches(0.5) if i % 2 == 0 else Inches(6.8)
        top = Inches(1.5) + (i // 2) * Inches(1.4)
        
        add_shape(slide9, left, top, Inches(5.8), Inches(1.2), LIGHT_GRAY)
        add_text_box(slide9, left + Inches(0.2), top + Inches(0.1), Inches(5.4), Inches(0.5),
                    module, font_size=16, bold=True, color=PRIMARY_BLUE)
        add_text_box(slide9, left + Inches(0.2), top + Inches(0.6), Inches(5.4), Inches(0.5),
                    pages, font_size=13, color=DARK_GRAY)
    
    add_slide_number(slide9, 9)
    
    # ==================== SLIDE 10: ADVANTAGES & LIMITATIONS ====================
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide10, WHITE)
    
    # Header bar
    add_shape(slide10, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide10, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Advantages & Limitations", font_size=36, bold=True, color=WHITE)
    
    # Advantages
    add_shape(slide10, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), RGBColor(212, 237, 218))
    add_text_box(slide10, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.5),
                 "✓ Advantages", font_size=22, bold=True, color=SUCCESS_GREEN)
    
    advantages = [
        "No mobile app installation required for students",
        "Real-time attendance tracking and updates",
        "Prevents proxy attendance with secure tokens",
        "Automatic absence marking saves faculty time",
        "Multiple export formats (PDF, Excel, CSV)",
        "Comprehensive analytics and insights",
        "Scalable architecture for future expansion",
        "Professional UI with responsive design"
    ]
    add_bullet_points(slide10, Inches(0.7), Inches(2.3), Inches(5.4), Inches(4.5),
                     advantages, font_size=14, color=DARK_GRAY)
    
    # Limitations
    add_shape(slide10, Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.5), RGBColor(248, 215, 184))
    add_text_box(slide10, Inches(7), Inches(1.6), Inches(5.7), Inches(0.5),
                 "⚠ Limitations", font_size=22, bold=True, color=RGBColor(255, 133, 27))
    
    limitations = [
        "Requires active internet connection",
        "QR code spoofing possible via screenshots",
        "Limited to classroom-based lectures",
        "Initial setup requires data entry",
        "No GPS-based location verification",
        "Dependency on device camera quality"
    ]
    add_bullet_points(slide10, Inches(7), Inches(2.3), Inches(5.7), Inches(4.5),
                     limitations, font_size=14, color=DARK_GRAY)
    
    add_slide_number(slide10, 10)
    
    # ==================== SLIDE 11: FUTURE ENHANCEMENTS ====================
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide11, WHITE)
    
    # Header bar
    add_shape(slide11, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_BLUE)
    add_text_box(slide11, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                 "Future Enhancements", font_size=36, bold=True, color=WHITE)
    
    # Enhancements
    enhancements = [
        ("Mobile App", "Native Android/iOS app for students with push notifications"),
        ("GPS Verification", "Location-based verification to prevent off-campus attendance"),
        ("Biometric Integration", "Fingerprint/Face ID for additional security"),
        ("AI Analytics", "Predictive analytics for at-risk students"),
        ("Parent Portal", "Real-time attendance notifications to parents"),
        ("Multi-Campus", "Support for multiple campuses and branches"),
        ("Offline Mode", "Queue attendance when offline, sync when connected"),
        ("REST API", "Full API for third-party integrations")
    ]
    
    for i, (title, desc) in enumerate(enhancements):
        left = Inches(0.5) if i % 2 == 0 else Inches(6.8)
        top = Inches(1.5) + (i // 2) * Inches(1.4)
        
        add_shape(slide11, left, top, Inches(5.8), Inches(1.2), LIGHT_GRAY)
        add_text_box(slide11, left + Inches(0.2), top + Inches(0.1), Inches(5.4), Inches(0.5),
                    title, font_size=16, bold=True, color=PRIMARY_BLUE)
        add_text_box(slide11, left + Inches(0.2), top + Inches(0.6), Inches(5.4), Inches(0.5),
                    desc, font_size=13, color=DARK_GRAY)
    
    add_slide_number(slide11, 11)
    
    # ==================== SLIDE 12: CONCLUSION ====================
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide12, PRIMARY_BLUE)
    
    # Title
    add_text_box(slide12, Inches(0.5), Inches(1), Inches(12), Inches(1),
                 "Conclusion", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Conclusion text
    conclusion = """AttendX successfully addresses the challenges of traditional attendance management 
by providing a modern, secure, and efficient QR code-based solution. The system demonstrates 
professional software engineering practices including clean architecture, security best practices, 
and comprehensive documentation."""
    add_text_box(slide12, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                 conclusion, font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Key achievements
    add_text_box(slide12, Inches(1), Inches(4), Inches(11), Inches(0.6),
                 "Key Achievements:", font_size=20, bold=True, color=WARNING_ORANGE, alignment=PP_ALIGN.CENTER)
    
    achievements = [
        "Contactless attendance marking without mobile apps",
        "Real-time tracking with live dashboard",
        "Secure QR codes with anti-replay protection",
        "Professional, portfolio-quality application"
    ]
    add_bullet_points(slide12, Inches(2), Inches(4.7), Inches(9), Inches(2.5),
                     achievements, font_size=16, color=WHITE)
    
    # Thank you
    add_text_box(slide12, Inches(0.5), Inches(6.5), Inches(12), Inches(0.8),
                 "Thank You!", font_size=28, bold=True, color=WARNING_ORANGE, alignment=PP_ALIGN.CENTER)
    
    add_slide_number(slide12, 12)
    
    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "AttendX_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    
    return output_path


if __name__ == "__main__":
    create_attendx_presentation()
